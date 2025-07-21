from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Recreate presentation after kernel reset
prs = Presentation()

# Define slide content again
slides_content = [
    ("Healthcare-Focused Identity Management System",
     "Name: Sowmya Venkat\nID: [Your ID]\nCourse: WILP Final Dissertation\nInstitution: BITS Pilani\nMentor: [Mentor's Name]\nDate: April 2025"),

    ("Introduction",
     "Identity and Access Management (IAM) is central to modern healthcare systems.\nProtecting sensitive patient data requires strict access policies.\nTraditional IAM solutions fall short in dynamic, multi-role healthcare environments.\nOur project proposes a secure, decentralized, AI-powered IAM system designed specifically for healthcare."),

    ("Problem Statement",
     "Current healthcare IAM systems are:\n- Centrally managed → single point of failure\n- Lacking anomaly detection\n- Trusting by default → prone to insider threats\nNo integration of Zero-Trust or decentralized identity practices.\nAdmins need better monitoring tools to maintain data integrity."),

    ("Project Objectives",
     "• Design and implement a secure IAM prototype for healthcare\n• Enforce Zero Trust principles at every access point\n• Use Blockchain for Decentralized Identity Management (DIM)\n• Enable intelligent anomaly detection using AI\n• Build an admin dashboard to monitor logs, anomalies, and users"),

    ("Project Scope",
     "User Roles: Admin, Doctor, Nurse, Patient, Guardian, Staff, Researcher\nSystem Modules:\n- Authentication & Role-Based Access\n- Blockchain-based logging\n- AI Anomaly Detection for login behavior\n- Real-time Monitoring via Dashboard"),

    ("System Architecture",
     "Components:\nFrontend: React.js\nBackend: Express.js (Node.js)\nDatabase: MongoDB\nBlockchain: Simulated ledger-based logging\nAI: Python (Scikit-learn, Isolation Forest)\n\n[Insert Architecture Diagram]"),

    ("Modules Implemented",
     "• Zero-Trust Identity Verification\n• RBAC/ABAC-Based Access Control\n• Blockchain-Powered DIM\n• AI-Powered Anomaly Detection\n• Admin Monitoring Dashboard"),

    ("Zero Trust Identity Verification",
     "No implicit trust—every access is authenticated\nJWT tokens issued with session expiry\nMiddleware protects all sensitive routes\nAll login attempts logged in DB and .json\n[Insert Screenshot: Login UI + Postman token]"),

    ("RBAC/ABAC-Based Access Control",
     "Access granted based on user role and attributes\nFrontend routing and backend middleware enforce security\nReal-time control over who accesses what\n[Insert Screenshot: Role-based dashboard UI]"),

    ("Blockchain-based DIM",
     "Sensitive operations logged to local blockchain\nEach block stores: action, userId, timestamp, hash\nMimics Hyperledger Fabric features\n[Insert Screenshot: AdminBlockchainLogs.jsx + ledger view]"),

    ("AI-Powered Anomaly Detection",
     "Every login saved in login_logs.json\nPython script detects outliers using Isolation Forest\nAdmins view anomalies via UI or backend\n[Insert Screenshot: Terminal + AdminAnalytics.jsx]"),

    ("Admin Dashboard",
     "React-based UI for monitoring\nIncludes: Manage Users, View Blockchain, Analytics, Logs\nProtected via route guards and JWT\n[Insert Screenshot: Admin Dashboard]"),

    ("Tech Stack Used",
     "Frontend: React, React-Bootstrap\nBackend: Node.js, Express.js\nDatabase: MongoDB\nBlockchain (Simulated): Node.js\nAI: Python, pandas, scikit-learn\nTools: Postman, Git, VS Code"),

    ("Results & Highlights",
     "Complete working prototype\nBlockchain ensures tamper-proof tracking\nAI model detects suspicious behavior\nLightweight, modular system ready for real-world deployment"),

    ("Conclusion & Future Work",
     "System meets objectives: security, decentralization, intelligence\nImproves healthcare IAM reliability and monitoring\nFuture Enhancements:\n- Real-time anomaly alerting\n- Hyperledger Fabric integration\n- MFA and geo-based access validation"),

    ("Thank You", "Questions?\nHappy to answer!")
]

# Function to add slides
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Add each slide
for title, content in slides_content:
    add_slide(title, content)

# Save presentation
pptx_path = "/mnt/data/Healthcare_IAM_Presentation.pptx"
prs.save(pptx_path)
pptx_path
